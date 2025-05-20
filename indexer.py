from typing import List, Dict
from datetime import datetime
import os
from docx import Document
from pptx import Presentation
import faiss
import numpy as np
from openai import OpenAI

# Globals for MVP
VECTOR_STORE_PATH = "vector.index"
CHUNK_SIZE = 500  # characters
EMBEDDING_DIM = 1536  # OpenAI text-embedding-ada-002
EMBED_MODEL = "text-embedding-ada-002"

# In-memory store for file metadata and chunks
indexed_files: List[Dict] = []
chunks: List[Dict] = []

def extract_docx_text(path: str) -> str:
    """Extract text from a .docx file."""
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_pptx_text(path: str) -> str:
    """Extract text from a .pptx file."""
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = getattr(shape, "text", "")
                if text and text.strip():
                    texts.append(text)
    return "\n".join(texts)

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    """Split text into chunks of a given size."""
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def get_embeddings(texts: List[str]) -> np.ndarray:
    """Get OpenAI embeddings for a list of texts."""
    client = OpenAI()
    response = client.embeddings.create(input=texts, model=EMBED_MODEL)
    return np.array([e.embedding for e in response.data], dtype=np.float32)

def index_documents() -> List[str]:
    """Index all .pptx and .docx files in the data folders."""
    global indexed_files, chunks
    indexed_files = []
    chunks = []
    all_texts = []
    chunk_meta = []
    folders = ["data/qdd", "data/qcp", "data/plr", "data/opsex"]
    for folder in folders:
        if not os.path.exists(folder):
            continue
        for fname in os.listdir(folder):
            if fname.endswith((".pptx", ".docx")):
                path = os.path.join(folder, fname)
                mtime = datetime.fromtimestamp(os.path.getmtime(path)).isoformat()
                indexed_files.append({"filename": path, "last_modified": mtime})
                # Extract text
                text = extract_docx_text(path) if fname.endswith(".docx") else extract_pptx_text(path)
                # Chunk and collect
                for i, chunk in enumerate(chunk_text(text)):
                    all_texts.append(chunk)
                    chunk_meta.append({"text": chunk, "file": path, "chunk_id": i})
    # Embed and build FAISS index
    if all_texts:
        embs = get_embeddings(all_texts)
        if embs.ndim == 1:
            embs = embs.reshape(1, -1)
        index = faiss.IndexFlatL2(EMBEDDING_DIM)
        index.add(embs)  # type: ignore
        faiss.write_index(index, VECTOR_STORE_PATH)
        chunks = chunk_meta
    return [f["filename"] for f in indexed_files]

def get_indexed_files() -> List[Dict]:
    """Return metadata for all indexed files."""
    return indexed_files
